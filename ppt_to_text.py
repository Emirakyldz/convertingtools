import os
import sys
from pathlib import Path
from pptx import Presentation

def pptx_to_text(pptx_path, txt_path):
    try:
        prs = Presentation(pptx_path)
        text = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + '\n'
        with open(txt_path, 'w', encoding='utf-8') as out:
            out.write(text)
        print(f"Converted: {pptx_path} -> {txt_path}")
    except Exception as e:
        print(f"Error occurred: {pptx_path}: {e}")

def main(pptx_folder):
    pptx_folder = Path(pptx_folder)
    if not pptx_folder.exists() or not pptx_folder.is_dir():
        print("Please enter a valid folder path.")
        return
    pptx_files = list(pptx_folder.glob('*.pptx'))
    if not pptx_files:
        print("No PPTX files found in the folder.")
        return
    for pptx_file in pptx_files:
        txt_file = pptx_file.with_suffix('.txt')
        pptx_to_text(pptx_file, txt_file)

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python ppt_to_text.py <pptx_folder>")
    else:
        main(sys.argv[1]) 